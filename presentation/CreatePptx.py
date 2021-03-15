import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Pt


def CreatePptx(car_specs,img_path, save_path, prices, home, page):
    prs = Presentation(home + "/Dropbox/DipMax/Code/Example_fb_content_page_v1.02.pptx")
    layout = prs.slide_layouts.get_by_name("Standard_fb_1.0")
    new_slide = prs.slides.add_slide(layout)

    # Insert name
    placeholder = new_slide.placeholders[0]
    placeholder.text = car_specs["Name"]

    # Insert category
    placeholder = new_slide.placeholders[15]
    placeholder.text = car_specs["Category"]

    # Insert the picture
    placeholder = new_slide.placeholders[13]
    picture = placeholder.insert_picture(img_path)

    # Insert Price
    placeholder = new_slide.placeholders[1]
    placeholder.text = prices[0] + "*\n" + prices[1] + "*\n" + prices[2] + "*"
    placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    placeholder.text_frame.paragraphs[1].alignment = PP_ALIGN.RIGHT
    placeholder.text_frame.paragraphs[2].alignment = PP_ALIGN.RIGHT

    # Add the table
    rows = 11
    cols = 2
    left = Cm(20.44)
    top = Cm(3.55)
    width = Cm(12.96)
    height = Cm(11.17)

    shapes = new_slide.shapes
    table = shapes.add_table(rows, cols, left, top, width, height)

    ## Define the style of table
    tbl = table._element.graphic.graphicData.tbl
    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id

    table = table.table
    table.columns[0].width = Cm(6.8)
    table.columns[1].width = Cm(6.16)
    table.first_row = False

    ##Hinzufügen des Tabelleninhaltes
    if page == "m":
        k = 0
        for key, value in car_specs.items():
            if key == "Name" or key == "Category" or key == "link" or key == "price:" or key == "net_price:":
                continue
            table.cell(k, 0).text_frame.text = key
            table.cell(k, 1).text = value
            if k == 10:
                break
            k += 1
    else:
        k = 0
        table.cell(0, 0).text_frame.text = "Range"
        table.cell(1, 0).text_frame.text = "City:"
        table.cell(1, 1).text = car_specs["City"]
        table.cell(2, 0).text_frame.text = "Highway:"
        table.cell(2, 1).text = car_specs["Highway"]
        table.cell(3, 0).text_frame.text = "Acceleration:"
        table.cell(3, 1).text = car_specs["Acceleration"]
        table.cell(4, 0).text_frame.text = "Top speed:"
        table.cell(4, 1).text = car_specs["Top Speed"]
        table.cell(5, 0).text_frame.text = "Number of Seats:"
        table.cell(5, 1).text = car_specs["Number of Seats"]
        table.cell(6, 0).text_frame.text = "Charge Time"
        table.cell(7, 0).text_frame.text = "Wall Plug:"
        table.cell(7, 1).text = car_specs["Wall Plug"]
        table.cell(8, 0).text_frame.text = "Wall Box:"
        table.cell(8, 1).text = car_specs["Wall box"]
        table.cell(9, 0).text_frame.text = "Station:"
        table.cell(9, 1).text = car_specs["Station"]
        table.cell(10, 0).text_frame.text = "Full Charge:"
        table.cell(10, 1).text = car_specs["Full charge"]

        ## Formatierung der Tabelle
    for i in range(0, cols):
        for k in range(0, rows):
            table.cell(k, i).margin_left = 0
            table.cell(k, i).margin_right = 0
            table.cell(k, i).margin_bottom = 0
            table.cell(k, i).margin_top = 0
            table.cell(k, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            table.cell(k, i).text_frame.paragraphs[0].font.size = Pt(20)
            table.cell(k, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            table.cell(k, i).fill.background()

    pptx_output = os.path.join(save_path, "fb_promotion_" + car_specs["Name"] + ".pptx")

    prs.save(pptx_output)

    return pptx_output